"""ドキュメントパーサー: 各種ファイル形式からテキストを抽出する"""

from __future__ import annotations

import csv
import hashlib
import io
import json as json_lib
from dataclasses import dataclass, field
from pathlib import Path

# ---------------------------------------------------------------------------
# 対応拡張子一覧（app.py側のfile_uploaderと合わせること）
# ---------------------------------------------------------------------------
SUPPORTED_EXTENSIONS = [
    ".pdf",
    ".xlsx", ".xls", ".ods",       # スプレッドシート系
    ".csv", ".tsv",                 # テキスト表形式
    ".docx", ".doc", ".rtf",       # ワープロ系
    ".txt", ".text", ".md",        # プレーンテキスト / Markdown
    ".json",                        # JSON
    ".html", ".htm",               # HTML
    ".pptx",                        # PowerPoint
]


@dataclass
class ParsedDocument:
    filename: str
    doc_type: str
    text: str
    page_count: int = 0
    sheet_names: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def parse_file(file_name: str, file_bytes: bytes) -> ParsedDocument:
    """ファイル拡張子に応じて適切なパーサーを呼び出す"""
    ext = Path(file_name).suffix.lower()

    dispatch = {
        ".pdf": _parse_pdf,
        ".xlsx": _parse_excel,
        ".xls": _parse_excel,
        ".ods": _parse_ods,
        ".csv": _parse_csv,
        ".tsv": _parse_tsv,
        ".docx": _parse_word,
        ".doc": _parse_doc_legacy,
        ".rtf": _parse_rtf,
        ".txt": _parse_text,
        ".text": _parse_text,
        ".md": _parse_text,
        ".json": _parse_json,
        ".html": _parse_html,
        ".htm": _parse_html,
        ".pptx": _parse_pptx,
    }

    parser = dispatch.get(ext)
    if parser:
        return parser(file_name, file_bytes)

    return ParsedDocument(
        filename=file_name,
        doc_type="Unknown",
        text="",
        warnings=[f"未対応の形式です: {ext}"],
    )


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def _parse_pdf(name: str, data: bytes) -> ParsedDocument:
    import pdfplumber

    warnings: list[str] = []
    pages_text: list[str] = []

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                tables = page.extract_tables()
                table_text = ""
                for table in tables:
                    for row in table:
                        cells = [str(c) if c else "" for c in row]
                        table_text += " | ".join(cells) + "\n"
                combined = text
                if table_text:
                    combined += "\n[表データ]\n" + table_text
                pages_text.append(combined)

            if not any(t.strip() for t in pages_text):
                warnings.append("テキスト抽出ができませんでした（スキャン画像の可能性）")

            return ParsedDocument(
                filename=name,
                doc_type="PDF",
                text="\n\n---\n\n".join(pages_text),
                page_count=len(pdf.pages),
                warnings=warnings,
            )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="PDF", text="",
            warnings=[f"PDF解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# Excel (.xlsx / .xls)
# ---------------------------------------------------------------------------
def _parse_excel(name: str, data: bytes) -> ParsedDocument:
    from openpyxl import load_workbook

    warnings: list[str] = []
    sheets_text: list[str] = []
    sheet_names: list[str] = []

    try:
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        for ws_name in wb.sheetnames:
            sheet_names.append(ws_name)
            ws = wb[ws_name]
            rows: list[str] = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_count += 1
                if row_count > 5000:
                    warnings.append(f"シート '{ws_name}' は5000行で切り捨てました")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            sheets_text.append(f"【シート: {ws_name}】\n" + "\n".join(rows))
        wb.close()

        return ParsedDocument(
            filename=name,
            doc_type="Excel",
            text="\n\n".join(sheets_text),
            sheet_names=sheet_names,
            warnings=warnings,
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="Excel", text="",
            warnings=[f"Excel解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# ODS (LibreOffice / Google Sheets export)
# ---------------------------------------------------------------------------
def _parse_ods(name: str, data: bytes) -> ParsedDocument:
    """ODS形式をpandasで読み込む。odfpy が必要。"""
    try:
        import pandas as pd

        sheets = pd.read_excel(io.BytesIO(data), engine="odf", sheet_name=None)
        parts: list[str] = []
        sheet_names: list[str] = []
        for sname, df in sheets.items():
            sheet_names.append(str(sname))
            text = df.to_csv(sep="|", index=False)
            parts.append(f"【シート: {sname}】\n{text}")

        return ParsedDocument(
            filename=name,
            doc_type="ODS",
            text="\n\n".join(parts),
            sheet_names=sheet_names,
        )
    except ImportError:
        return ParsedDocument(
            filename=name, doc_type="ODS", text="",
            warnings=["ODS読み込みには odfpy パッケージが必要です (pip install odfpy)"],
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="ODS", text="",
            warnings=[f"ODS解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------
def _parse_csv(name: str, data: bytes) -> ParsedDocument:
    try:
        text = _decode_bytes(data)
        reader = csv.reader(io.StringIO(text))
        rows = [" | ".join(row) for row in reader if any(row)]
        return ParsedDocument(
            filename=name, doc_type="CSV",
            text="\n".join(rows[:5000]),
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="CSV", text="",
            warnings=[f"CSV解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# TSV
# ---------------------------------------------------------------------------
def _parse_tsv(name: str, data: bytes) -> ParsedDocument:
    try:
        text = _decode_bytes(data)
        reader = csv.reader(io.StringIO(text), delimiter="\t")
        rows = [" | ".join(row) for row in reader if any(row)]
        return ParsedDocument(
            filename=name, doc_type="TSV",
            text="\n".join(rows[:5000]),
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="TSV", text="",
            warnings=[f"TSV解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------
def _parse_word(name: str, data: bytes) -> ParsedDocument:
    from docx import Document

    warnings: list[str] = []
    try:
        doc = Document(io.BytesIO(data))
        paragraphs: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                prefix = ""
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    prefix = "## "
                paragraphs.append(prefix + para.text)

        for table in doc.tables:
            trows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                trows.append(" | ".join(cells))
            if trows:
                paragraphs.append("[表データ]\n" + "\n".join(trows))

        return ParsedDocument(
            filename=name,
            doc_type="Word",
            text="\n".join(paragraphs),
            page_count=len(doc.paragraphs) // 30 or 1,
            warnings=warnings,
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="Word", text="",
            warnings=[f"Word解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# Word (.doc) - レガシー形式
# ---------------------------------------------------------------------------
def _parse_doc_legacy(name: str, data: bytes) -> ParsedDocument:
    """古い.doc形式。テキスト部分のバイナリ抽出を試みる。"""
    warnings: list[str] = []
    try:
        # バイナリからテキスト部分を抽出（簡易的）
        text = data.decode("utf-8", errors="ignore")
        # 制御文字を除去
        cleaned = "".join(c for c in text if c.isprintable() or c in "\n\r\t")
        lines = [l.strip() for l in cleaned.splitlines() if l.strip() and len(l.strip()) > 3]
        if not lines:
            warnings.append(".doc形式のテキスト抽出は限定的です。.docx形式への変換を推奨します")
        return ParsedDocument(
            filename=name,
            doc_type="Word (legacy)",
            text="\n".join(lines[:3000]),
            warnings=warnings,
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="Word (legacy)", text="",
            warnings=[f".doc解析エラー: {e}。.docx形式への変換を推奨します"],
        )


# ---------------------------------------------------------------------------
# RTF
# ---------------------------------------------------------------------------
def _parse_rtf(name: str, data: bytes) -> ParsedDocument:
    """RTF形式。striprtf パッケージがあれば使用、なければ簡易抽出。"""
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(data.decode("utf-8", errors="ignore"))
        return ParsedDocument(filename=name, doc_type="RTF", text=text)
    except ImportError:
        # striprtf がない場合は簡易抽出
        raw = data.decode("utf-8", errors="ignore")
        import re
        text = re.sub(r"[\\{}\[\]]", " ", raw)
        text = re.sub(r"\\[a-z]+\d*\s?", "", text)
        cleaned = " ".join(text.split())
        return ParsedDocument(
            filename=name, doc_type="RTF", text=cleaned,
            warnings=["RTFの高精度解析には striprtf パッケージが必要です (pip install striprtf)"],
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="RTF", text="",
            warnings=[f"RTF解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# プレーンテキスト / Markdown
# ---------------------------------------------------------------------------
def _parse_text(name: str, data: bytes) -> ParsedDocument:
    ext = Path(name).suffix.lower()
    doc_type = "Markdown" if ext == ".md" else "テキスト"
    try:
        text = _decode_bytes(data)
        return ParsedDocument(filename=name, doc_type=doc_type, text=text[:100000])
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type=doc_type, text="",
            warnings=[f"テキスト解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# JSON
# ---------------------------------------------------------------------------
def _parse_json(name: str, data: bytes) -> ParsedDocument:
    try:
        text = _decode_bytes(data)
        obj = json_lib.loads(text)
        pretty = json_lib.dumps(obj, ensure_ascii=False, indent=2)
        return ParsedDocument(
            filename=name, doc_type="JSON",
            text=pretty[:100000],
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="JSON", text="",
            warnings=[f"JSON解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------
def _parse_html(name: str, data: bytes) -> ParsedDocument:
    """HTMLからテキストを抽出。BeautifulSoup があれば使用。"""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(data, "html.parser")
        # スクリプトとスタイルを除去
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return ParsedDocument(filename=name, doc_type="HTML", text=text[:100000])
    except ImportError:
        # bs4 がない場合は簡易抽出
        import re
        raw = _decode_bytes(data)
        text = re.sub(r"<script[^>]*>.*?</script>", "", raw, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = " ".join(text.split())
        return ParsedDocument(
            filename=name, doc_type="HTML", text=text[:100000],
            warnings=["HTML高精度解析には beautifulsoup4 が必要です (pip install beautifulsoup4)"],
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="HTML", text="",
            warnings=[f"HTML解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------
def _parse_pptx(name: str, data: bytes) -> ParsedDocument:
    """PowerPointからテキストを抽出。python-pptx が必要。"""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                slides_text.append(f"【スライド {i}】\n" + "\n".join(texts))

        return ParsedDocument(
            filename=name,
            doc_type="PowerPoint",
            text="\n\n".join(slides_text),
            page_count=len(prs.slides),
        )
    except ImportError:
        return ParsedDocument(
            filename=name, doc_type="PowerPoint", text="",
            warnings=["PowerPoint読み込みには python-pptx が必要です (pip install python-pptx)"],
        )
    except Exception as e:
        return ParsedDocument(
            filename=name, doc_type="PowerPoint", text="",
            warnings=[f"PowerPoint解析エラー: {e}"],
        )


# ---------------------------------------------------------------------------
# ユーティリティ
# ---------------------------------------------------------------------------
def _decode_bytes(data: bytes) -> str:
    """バイト列を文字列にデコード。日本語エンコーディングを自動検出。"""
    for enc in ("utf-8", "utf-8-sig", "shift_jis", "cp932", "euc-jp", "iso-2022-jp"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


# ---------------------------------------------------------------------------
# 財務CSV自動抽出（LLMを経由せず直接構造化データとして抽出）
# ---------------------------------------------------------------------------
def _parse_num(s: str) -> int | None:
    """数値文字列をintに変換"""
    s = s.strip().replace(",", "")
    if not s or s == "-":
        return None
    try:
        return int(s)
    except ValueError:
        try:
            return int(float(s))
        except ValueError:
            return None


def _readable_yen(val: int) -> str:
    """金額を読みやすい日本語表記で返す"""
    abs_val = abs(val)
    if abs_val >= 100_000_000:
        return f"約{val / 100_000_000:.1f}億円"
    elif abs_val >= 10_000:
        return f"約{val / 10_000:.0f}万円"
    return f"{val:,}円"


def _find_row(
    rows: list[list[str]], patterns: list[str], cols: tuple[int, ...] = (0,)
) -> list[str] | None:
    """パターンに一致し数値データを含む行を返す"""
    for row in rows:
        if not row or len(row) < 4:
            continue
        for c in cols:
            if c >= len(row):
                continue
            label = row[c].strip()
            for p in patterns:
                if label == p:
                    for i in range(3, min(len(row), 20)):
                        if _parse_num(row[i]) is not None:
                            return row
    return None


_PL_ITEMS: list[tuple[str, list[str]]] = [
    ("revenue", ["売上高合計", "売上高計"]),
    ("cost_of_sales", ["売上原価合計", "売上原価計"]),
    ("gross_profit", ["売上総利益", "売上総利益合計"]),
    ("sga", ["販売費及び一般管理費合計", "販売費及一般管理費合計", "販管費合計"]),
    ("operating_profit", ["営業利益", "営業利益合計"]),
    ("ordinary_profit", ["経常利益", "経常利益合計"]),
    ("pretax_income", ["税引前当期純利益", "税引前当期純利益合計"]),
    ("net_income", ["当期純利益", "当期純利益合計"]),
]

_BS_ITEMS: list[tuple[str, list[str], tuple[int, ...]]] = [
    ("total_assets", ["資産の部合計"], (0,)),
    ("current_assets", ["流動資産合計"], (0,)),
    ("cash", ["現金及び預金合計", "現金及び預金", "現金及預金合計"], (0, 1)),
    ("fixed_assets", ["固定資産合計"], (0,)),
    ("total_liabilities", ["負債の部合計"], (0,)),
    ("current_liabilities", ["流動負債合計"], (0,)),
    ("fixed_liabilities", ["固定負債合計"], (0,)),
    ("interest_bearing_debt", ["長期借入金合計", "長期借入金"], (0, 1)),
    ("net_assets", ["純資産の部合計"], (0,)),
]

# P/L項目の表示名マッピング
PL_DISPLAY_NAMES = {
    "revenue": "売上高",
    "cost_of_sales": "売上原価",
    "gross_profit": "売上総利益",
    "sga": "販管費",
    "operating_profit": "営業利益",
    "ordinary_profit": "経常利益",
    "pretax_income": "税引前利益",
    "net_income": "当期純利益",
}

BS_DISPLAY_NAMES = {
    "total_assets": "総資産",
    "current_assets": "流動資産",
    "cash": "現預金",
    "fixed_assets": "固定資産",
    "total_liabilities": "負債合計",
    "current_liabilities": "流動負債",
    "fixed_liabilities": "固定負債",
    "interest_bearing_debt": "有利子負債",
    "net_assets": "純資産",
}


@dataclass
class ExtractedPL:
    """CSVから直接抽出したP/Lデータ"""
    filename: str
    fiscal_month: str  # e.g., "8"
    period_label: str = ""  # LLMが特定した年度ラベル or 自動生成
    revenue: int | None = None
    cost_of_sales: int | None = None
    gross_profit: int | None = None
    sga: int | None = None
    operating_profit: int | None = None
    ordinary_profit: int | None = None
    pretax_income: int | None = None
    net_income: int | None = None


@dataclass
class ExtractedBS:
    """CSVから直接抽出したB/Sデータ"""
    filename: str
    fiscal_month: str
    period_label: str = ""
    total_assets: int | None = None
    current_assets: int | None = None
    cash: int | None = None
    fixed_assets: int | None = None
    total_liabilities: int | None = None
    current_liabilities: int | None = None
    fixed_liabilities: int | None = None
    interest_bearing_debt: int | None = None
    net_assets: int | None = None


@dataclass
class ExtractedFinancials:
    """CSVから直接抽出した全財務データ"""
    pl_list: list[ExtractedPL] = field(default_factory=list)
    bs_list: list[ExtractedBS] = field(default_factory=list)

    def to_dict(self) -> dict:
        """JSON保存用に辞書化"""
        return {
            "pl_list": [
                {k: v for k, v in pl.__dict__.items()}
                for pl in self.pl_list
            ],
            "bs_list": [
                {k: v for k, v in bs.__dict__.items()}
                for bs in self.bs_list
            ],
        }

    @staticmethod
    def from_dict(d: dict) -> "ExtractedFinancials":
        """辞書から復元"""
        if not d:
            return ExtractedFinancials()
        return ExtractedFinancials(
            pl_list=[ExtractedPL(**p) for p in d.get("pl_list", [])],
            bs_list=[ExtractedBS(**b) for b in d.get("bs_list", [])],
        )


def _extract_pl_from_csv(fname: str, text: str) -> ExtractedPL | None:
    """損益計算書CSVから構造化P/Lデータを抽出"""
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if len(rows) < 3:
        return None

    header = rows[0]
    month_cols: list[tuple[int, str]] = []
    total_col: int | None = None
    for i in range(3, len(header)):
        h = header[i].strip()
        if h == "合計":
            total_col = i
        elif h and h != "決算整理":
            month_cols.append((i, h))

    if total_col is None:
        return None

    # 決算月を検出
    fiscal_month = ""
    if month_cols:
        fiscal_month = month_cols[-1][1].replace("月", "").strip()

    pl = ExtractedPL(filename=fname, fiscal_month=fiscal_month)

    found_any = False
    for key, patterns in _PL_ITEMS:
        row = _find_row(rows, patterns, cols=(0,))
        if row and total_col < len(row):
            val = _parse_num(row[total_col])
            if val is not None:
                setattr(pl, key, val)
                found_any = True

    return pl if found_any else None


def _extract_bs_from_csv(fname: str, text: str) -> ExtractedBS | None:
    """貸借対照表CSVから構造化B/Sデータを抽出"""
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if len(rows) < 3:
        return None

    header = rows[0]
    month_cols: list[tuple[int, str]] = []
    for i in range(3, len(header)):
        h = header[i].strip()
        if h:
            month_cols.append((i, h))

    if not month_cols:
        return None

    # 期末（最終月の列）を使用
    latest_col = month_cols[-1][0]
    fiscal_month = month_cols[-1][1].replace("月", "").strip()

    bs = ExtractedBS(filename=fname, fiscal_month=fiscal_month)

    found_any = False
    for key, patterns, cols in _BS_ITEMS:
        row = _find_row(rows, patterns, cols=cols)
        if row and latest_col < len(row):
            val = _parse_num(row[latest_col])
            if val is not None:
                setattr(bs, key, val)
                found_any = True

    return bs if found_any else None


def _dedup_files(files: list[tuple[str, bytes]]) -> list[tuple[str, bytes]]:
    """同一内容のファイルを除外（ファイル名ではなく内容のMD5ハッシュで判定）"""
    seen: set[str] = set()
    result: list[tuple[str, bytes]] = []
    for fname, data in files:
        h = hashlib.md5(data).hexdigest()
        if h not in seen:
            seen.add(h)
            result.append((fname, data))
    return result


def extract_financial_data(files: list[tuple[str, bytes]]) -> ExtractedFinancials:
    """CSVから財務データを直接構造化データとして抽出する（LLM不要）

    P/Lは売上高昇順（=時系列順と推定）でソートして返す。
    B/Sは総資産昇順でソートして返す。
    """
    unique_files = _dedup_files(files)

    pl_list: list[ExtractedPL] = []
    bs_list: list[ExtractedBS] = []

    for fname, data in unique_files:
        if not fname.lower().endswith(".csv"):
            continue
        try:
            text = _decode_bytes(data)
        except Exception:
            continue

        if "損益" in fname:
            pl = _extract_pl_from_csv(fname, text)
            if pl:
                pl_list.append(pl)
        elif "貸借" in fname:
            bs = _extract_bs_from_csv(fname, text)
            if bs:
                bs_list.append(bs)

    # 売上高昇順でソート（= 時系列順と推定）
    pl_list.sort(key=lambda p: p.revenue or 0)
    # 総資産昇順でソート
    bs_list.sort(key=lambda b: b.total_assets or 0)

    # デフォルトの期間ラベルを生成（売上高ベース）
    for i, pl in enumerate(pl_list):
        rev_str = _readable_yen(pl.revenue) if pl.revenue else "不明"
        month_str = f"{pl.fiscal_month}月期" if pl.fiscal_month else ""
        pl.period_label = f"第{i+1}期 ({month_str} 売上{rev_str})"

    for i, bs in enumerate(bs_list):
        asset_str = _readable_yen(bs.total_assets) if bs.total_assets else "不明"
        month_str = f"{bs.fiscal_month}月期" if bs.fiscal_month else ""
        bs.period_label = f"第{i+1}期 ({month_str} 総資産{asset_str})"

    return ExtractedFinancials(pl_list=pl_list, bs_list=bs_list)


def apply_period_labels_from_llm(
    financials: ExtractedFinancials,
    llm_pl_trends: list[dict],
    llm_bs_trends: list[dict],
) -> None:
    """LLMが特定した年度ラベルをCSV抽出データにマッチングして適用する。

    マッチングは売上高/総資産の近似値で行う（±10%以内でマッチ）。
    マッチしない場合はデフォルトラベルを維持。
    """
    import re

    def _is_valid_period(label: str) -> bool:
        """「XXXX年X月期」形式かチェック"""
        return bool(re.match(r"^\d{4}年\d{1,2}月期$", label.strip()))

    # P/Lマッチング
    for llm_pl in llm_pl_trends:
        period = llm_pl.get("period", "")
        llm_rev = llm_pl.get("revenue")
        if not _is_valid_period(period) or llm_rev is None:
            continue
        for pl in financials.pl_list:
            if pl.revenue is None:
                continue
            # ±10%以内でマッチ
            if abs(pl.revenue - llm_rev) / max(pl.revenue, 1) < 0.10:
                pl.period_label = period
                break

    # B/Sマッチング
    for llm_bs in llm_bs_trends:
        period = llm_bs.get("period", "")
        llm_assets = llm_bs.get("total_assets")
        if not _is_valid_period(period) or llm_assets is None:
            continue
        for bs in financials.bs_list:
            if bs.total_assets is None:
                continue
            if abs(bs.total_assets - llm_assets) / max(bs.total_assets, 1) < 0.10:
                bs.period_label = period
                break


def extract_financial_summary(files: list[tuple[str, bytes]]) -> str:
    """財務CSVから主要データを自動抽出し、LLM用のテキストで返す。

    LLMには年度ラベル特定の参考情報として渡す。
    実際の数値表示にはextract_financial_data()の結果を使う。
    """
    unique_files = _dedup_files(files)

    pl_parts: list[str] = []
    bs_parts: list[str] = []

    for fname, data in unique_files:
        if not fname.lower().endswith(".csv"):
            continue
        try:
            text = _decode_bytes(data)
        except Exception:
            continue

        if "損益" in fname:
            pl = _extract_pl_from_csv(fname, text)
            if pl:
                rev_str = f"{pl.revenue:,}円（{_readable_yen(pl.revenue)}）" if pl.revenue else "不明"
                month_str = f"{pl.fiscal_month}月期" if pl.fiscal_month else ""
                lines = [
                    f"■ 損益計算書 [{month_str}決算] 売上高={rev_str}",
                    f"  ファイル: {fname}",
                    "  単位: 円",
                ]
                for key, _ in _PL_ITEMS:
                    val = getattr(pl, key, None)
                    if val is not None:
                        lines.append(f"  {PL_DISPLAY_NAMES[key]}: {val:,}円")
                pl_parts.append("\n".join(lines))

        elif "貸借" in fname:
            bs = _extract_bs_from_csv(fname, text)
            if bs:
                asset_str = f"{bs.total_assets:,}円" if bs.total_assets else "不明"
                month_str = f"{bs.fiscal_month}月期" if bs.fiscal_month else ""
                lines = [
                    f"■ 貸借対照表 [{month_str}] 総資産={asset_str}",
                    f"  ファイル: {fname}",
                    "  単位: 円",
                ]
                for key, _, _ in _BS_ITEMS:
                    val = getattr(bs, key, None)
                    if val is not None:
                        lines.append(f"  {BS_DISPLAY_NAMES[key]}: {val:,}円")
                bs_parts.append("\n".join(lines))

    # 売上高でソートして表示
    parts = sorted(pl_parts) + sorted(bs_parts)
    if not parts:
        return ""

    return (
        "【財務データ（CSVから自動抽出・正確な数値）】\n"
        f"P/L: {len(pl_parts)}期分、B/S: {len(bs_parts)}期分のデータを抽出済み。\n"
        "以下は各CSVの年間合計値です。\n"
        "IMの記載と照合して、各CSVが何年度のものか特定してください。\n"
        "pl_trends/bs_trendsのperiodには必ず「XXXX年X月期」形式で出力してください。\n"
        "「期間A」「CSV#1」等の仮ラベルは禁止です。\n\n"
        + "\n\n".join(parts)
    )
